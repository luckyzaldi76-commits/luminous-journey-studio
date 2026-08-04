import json
from pathlib import Path
from pptx import Presentation


def create_ppt(json_file):

    with open(json_file, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()

    layout = prs.slide_layouts[1]

    for key, value in data.items():

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = key.replace("_", " ").title()

        slide.placeholders[1].text = str(value)

    output_dir = Path("05_PPT")
    output_dir.mkdir(exist_ok=True)

    output_file = output_dir / (Path(json_file).stem + ".pptx")

    prs.save(output_file)

    return output_file