from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt


class PPTBuilder:

    def build(
        self,
        data,
        output_file
    ):

        prs = Presentation()

        def add_slide(title, content):

            slide = prs.slides.add_slide(
                prs.slide_layouts[5]
            )

            slide.shapes.title.text = title

            textbox = slide.shapes.add_textbox(
                Inches(0.8),
                Inches(1.2),
                Inches(8),
                Inches(5)
            )

            tf = textbox.text_frame

            p = tf.paragraphs[0]

            p.text = content

            p.font.size = Pt(20)

        for key, value in data.items():

            add_slide(
                key,
                value
            )

        prs.save(output_file)

        print(f"PPT Saved : {output_file}")


builder = PPTBuilder()